# utils/file_handler.py

import fitz  # PyMuPDF
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import markdownify

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    text = ""
    estrutura = []
    for page_num, page in enumerate(doc):
        content = page.get_text()
        text += content + "\n"
        estrutura.append({"pagina": page_num + 1, "conteudo": content.strip()})
    return text.strip(), estrutura

def extract_text_from_docx(path):
    doc = Document(path)
    text = ""
    estrutura = []
    for i, p in enumerate(doc.paragraphs):
        content = p.text.strip()
        if content:
            text += content + "\n"
            estrutura.append({"paragrafo": i + 1, "conteudo": content})
    return text.strip(), estrutura

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = ""
    estrutura = []
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        estrutura.append({"slide": i + 1, "conteudo": slide_text.strip()})
        text += slide_text + "\n"
    return text.strip(), estrutura

def extract_text_from_html(path):
    with open(path, "r", encoding="utf-8") as file:
        soup = BeautifulSoup(file, "html.parser")
        text = soup.get_text(separator="\n").strip()
    estrutura = [{"tipo": "html", "conteudo": text}]
    return text, estrutura

def convert_to_markdown(text):
    # Você pode adaptar isso depois pra converter melhor de acordo com o tipo
    return markdownify.markdownify(text, heading_style="ATX")
